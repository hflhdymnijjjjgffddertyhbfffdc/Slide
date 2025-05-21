#!/usr/bin/env python
# -*- coding: utf-8 -*-
#读取PPT备注
import os
import tempfile
from subprocess import call
import subprocess
from pathlib import Path

from pdf2image import convert_from_path
from pptx import Presentation
import requests
import hashlib

FFMPEG_NAME = 'ffmpeg'

def get_wav(text, file_name):
    url = "http://127.0.0.1:9880"  # 替换为你的实际 URL
    data = {
        "text": text,
        "text_language": "zh",
        "cut_punc": "，。",
        "speed": 1.5
    }

    response = requests.post(url, json=data)

    if response.status_code == 200:
        with open(file_name, 'wb') as f:
            f.write(response.content)
        print("音频已成功保存为" + file_name)
    elif response.status_code == 400:
        error_info = response.json()
        print("请求失败:", error_info)
    else:
        print(f"请求返回了意外的状态码: {response.status_code}")
        print(f"音频生成失败，文件名: {file_name}")
        open(file_name, 'a').close()  # 创建空文件防止后续流程中断

def generate_silence_wav(path, duration=1):
    call([
        'ffmpeg',
        '-f', 'lavfi',
        '-i', 'anullsrc=channel_layout=stereo:sample_rate=44100',
        '-t', str(duration),
        '-q:a', '9',
        '-acodec', 'libmp3lame',
        path
    ])

def calculate_hash(file_path):
    sha256_hash = hashlib.sha256()
    with open(file_path, "rb") as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()


def ppt_presenter(pptx_path, pdf_path, output_path):
    with tempfile.TemporaryDirectory() as temp_path:
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)
        
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            notes = ""

            # 优先读取备注
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            notes += text + "。"
            notes = notes.replace("。。", "。").replace("  ", " ").replace("-", " ").strip()

            # 如果备注为空，则读取正文文本框中的所有文本
            if not notes:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            texts.append(text)
                notes = "。".join(texts)

            print(f"Slide {i+1} 使用文本: {notes}")

            image_path = os.path.join(temp_path, f'frame_{i}.jpg')
            audio_path = os.path.join(temp_path, f'frame_{i}.mp3')
            image.save(image_path)

            if notes:
                get_wav(notes, audio_path)
            else:
                # 生成1秒静音音频，防止ffmpeg报错
                generate_silence_wav(audio_path)

            ffmpeg_call(image_path, audio_path, temp_path, i)

        video_list = [os.path.join(temp_path, f'frame_{i}.ts') for i in range(len(images_from_path))]
        missing_files = [file for file in video_list if not os.path.exists(file)]
        if missing_files:
            print(f"以下文件缺失，无法继续合并：{missing_files}")
            return

        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)


def ffmpeg_call(image_path, audio_path, temp_path, i):
    out_path_mp4 = os.path.join(temp_path, f'frame_{i}.mp4')
    out_path_ts = os.path.join(temp_path, f'frame_{i}.ts')
    
    call([
        "ffmpeg",
        "-loop", "1", "-y",
        "-i", image_path,
        "-i", audio_path,
        "-vf", "scale=trunc(iw/2)*2:ih",
        "-c:v", "libx264",
        "-tune", "stillimage",
        "-c:a", "aac",
        "-b:a", "192k",
        "-pix_fmt", "yuv420p",
        "-movflags", "+faststart",
        "-shortest", out_path_mp4
    ])

    call([FFMPEG_NAME, '-y', '-i', out_path_mp4, '-c', 'copy',
          '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts])

def ffmpeg_concat(video_list_str, out_path):
    call([FFMPEG_NAME, '-y', '-f', 'mpegts', '-i', f'{video_list_str}',
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])

def convert_pptx_to_pdf(input_file, output_file, dir_pdf):
    try:
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', dir_pdf, input_file], check=True)
        print(f"转换成功: {output_file}")
        return True
    except subprocess.CalledProcessError as e:
        print(f"转换失败: {e}")
        return False

def make_dir(folder):
    if not os.path.exists(folder):
        os.makedirs(folder)

def main():
    dir_pdf = "/home/asus/HuFan/ppt/pdf1"
    dir_ppt = "/home/asus/HuFan/ppt/ppt1"
    dir_mp4 = "/home/asus/HuFan/ppt/mp4"
    make_dir(dir_mp4)
    
    folder = Path(dir_ppt)
    ppt_files = folder.glob('*.pptx')
    
    for ppt_file in ppt_files:
        print("-"*30, ppt_file)
        
        pdf_path = os.path.join(dir_pdf, ppt_file.name.replace(".pptx", ".pdf"))
        mp4_path = os.path.join(dir_mp4, ppt_file.name.replace(".pptx", ".mp4"))
    
        convert_pptx_to_pdf(ppt_file, pdf_path, dir_pdf)
        ppt_presenter(ppt_file, pdf_path, mp4_path)

if __name__ == '__main__':
    main()
