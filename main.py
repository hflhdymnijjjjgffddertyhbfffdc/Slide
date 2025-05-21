import os
from pathlib import Path
import openai
import time
import json
import requests
import subprocess
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

import base64
import io
from pptx.dml.color import RGBColor



import logging



class ColoredFormatter(logging.Formatter):
    COLORS = {
        'DEBUG': '\033[94m',  # 蓝色
        'INFO': '\033[92m',   # 绿色
        'WARNING': '\033[93m', # 黄色
        'ERROR': '\033[91m',  # 红色
        'CRITICAL': '\033[95m', # 紫色
    }
    RESET = '\033[0m'

    def format(self, record):
        levelname = record.levelname
        color = self.COLORS.get(levelname, self.RESET)
        record.msg = f"{color}{record.msg}{self.RESET}"
        return super().format(record)

def setup_logger():
    logger = logging.getLogger()
    logger.setLevel(logging.INFO)  # 设置日志级别

    # 创建控制台处理器
    ch = logging.StreamHandler()
    ch.setLevel(logging.INFO)

    # 创建带颜色的格式化器
    formatter = ColoredFormatter('%(asctime)s - %(levelname)s - %(message)s')
    ch.setFormatter(formatter)

    # 添加处理器到logger
    logger.addHandler(ch)
    

setup_logger()



## =============================================================================== gpt
def ask_gpt(query, streaming_flg=True, max_retries=5, flag=True):
    if flag:
        openai.api_base = "https://opus.gptuu.com/v1"
        openai.api_key = "sk-G0TnsMYO17kGKpQn0ScLk6xvtL72iKCiFkM4CSGfhxNRIQR6"
    else:
        openai.api_base = "https://openkey.cloud/v1"
        openai.api_key = "sk-JhthldK6DOHIqh7HA29e61451116419aA96d526b22886604"
        
    for attempt in range(max_retries + 1):
        try:
            start_time = time.time()
            completion = openai.ChatCompletion.create(
                model="gpt-4o-mini",
                messages=[{'role': 'user', 'content': query}],
                stream=streaming_flg
            )
            msg = None

            if streaming_flg:
                msg, completion = process_streaming_response(completion)
            else:
                if not hasattr(completion, 'choices') or not completion.choices:
                    raise ValueError("No choices returned from API.")
                msg = completion.choices[0].message['content']

            # 如果成功获取到消息，返回结果
            if msg:
                end_time = time.time()
                logging.info(f"耗时: {end_time - start_time:.2f} 秒    GPT output: {msg}")
                time.sleep(3)
                return msg, completion

        except Exception as err:
            logging.warning(f'OpenAI API Error: {str(err)}')
            if attempt < max_retries:
                logging.warning(f"Retrying... (Attempt {attempt + 1} of {max_retries})")
                time.sleep(2 ** (attempt + 2))  # 指数退避
            else:
                if flag:
                    logging.error("Max retries reached. Change OpenAI.")
                    return ask_gpt(query, streaming_flg, flag=False)
                else:
                    logging.error("Max retries reached. Change OpenAI.")
                    return ask_gpt(query, streaming_flg, flag=True)
            



def process_streaming_response(completion):
    llist = []
    msg = ""
    for i in completion:
        if len(i['choices']) != 0:
            delta = i['choices'][0]['delta']
            content = delta.get('content', None)
            if content is not None:
                msg += content
        else:
            continue

    def generator(llist):
        for i in llist:
            yield i
    return msg, generator(llist)





## =============================================================================== ppt

def set_textbox_font_size(textbox, value):
    paragraphs = textbox.text_frame.paragraphs
    for paragraph in paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(value)
            

def set_textbox_align_center(textbox):
    paragraphs = textbox.text_frame.paragraphs
    for paragraph in paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        


def set_textbox_autowrap(shape):
    shape.text_frame.word_wrap = True
    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    

def new_slide(ppt, title, context, first=False):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    apply_slide_theme(slide)

    # 标题设置
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1.5))
    title_box.text = title
    set_title_style(title_box)
    
    # 内容区域设置
    content_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(2), width=Inches(18), height=Inches(6))
    content_box.text_frame.word_wrap = True
    
    # 添加段落并设置样式
    paragraphs = context.split("\n")
    for para in paragraphs:
        p = content_box.text_frame.add_paragraph()
        p.text = para.strip()  # 去除多余的空格
        p.alignment = PP_ALIGN.CENTER  # 段落居中
        p.level = 1  # 设置段落缩进级别
        p.font.size = Pt(24)  # 设置字体大小
        p.font.color.rgb = RGBColor(64, 64, 64)  # 设置字体颜色
        p.line_spacing = 1.2  # 设置行间距
    
    # 首页特殊处理
    if first:
        subtitle_box = slide.shapes.add_textbox(Inches(4), Inches(4), Inches(18), Inches(6))
        subtitle_box.text = context
        subtitle_box.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        subtitle_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_base64_pic(src, info, slide):
    """添加带说明的图片"""
    # 图片说明框
    caption_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    caption_box.text = info
    caption_box.fill.solid()
    caption_box.fill.fore_color.rgb = RGBColor(240, 240, 240)  # 浅灰背景
    caption_box.text_frame.paragraphs[0].font.size = Pt(18)
    caption_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)

    img_stream = io.BytesIO(base64.b64decode(src.split(',')[1]))
    pic = slide.shapes.add_picture(
        img_stream, 
        left=Inches(2), 
        top=Inches(1.5), 
        width=Inches(12)
    )
    
    # 添加图片边框
    line = pic.line
    line.color.rgb = RGBColor(200, 200, 200)
    line.width = Pt(1.5)






## =============================================================================== ppt
def apply_slide_theme(slide):
    """设置幻灯片背景和基础样式"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)  # 浅灰色背景

def set_title_style(textbox):
    """设置标题样式"""
    textbox.text_frame.paragraphs[0].font.size = Pt(40)
    textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 90, 158)  # 深蓝色
    textbox.text_frame.paragraphs[0].font.bold = True
    textbox.text_frame.paragraphs[0].space_after = Pt(15)  # 段后间距
    
def set_content_style(textbox):
    """设置正文样式"""
    textbox.text_frame.paragraphs[0].font.size = Pt(24)
    textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)  # 深灰色
    textbox.text_frame.paragraphs[0].line_spacing = 1.2  # 行间距    








    
# def generate_ppt(content, file_name, dir_ppt):
#     """生成研究论文PPT的完整函数
    
#     Args:
#         content (list): 包含文档内容和图片的列表
#         file_name (str): 原始文件名
#         dir_ppt (str): PPT输出目录
#     """
#     # ================== 初始化处理 ==================
#     # 提取图片数据
#     pics = ("[num1]:data:image" + content[1]).replace("\n\n", "\n").split("\n")

#     # ================== 第一阶段：结构化数据提取 ==================
#     struct_prompt = f"""
#     【严格JSON生成任务】
#     请从以下文档中提取：
#     - title（英文原题）
#     - author（作者列表）
#     - organization（单位列表，无则保留空字符串）
#     - abstract（中文摘要）
#     - background（研究背景）
#     - contribution（本文贡献）
#     - method（方法）
#     - experiment（实验）
#     - conclusion（结论）

#     输出要求：
#     1. 严格遵循JSON格式
#     2. 使用```json代码块包裹
#     3. 除标题作者单位外全用中文
#     4. 空值字段保持空字符串

#     文档内容：
#     {content[0]}
#     """
    
#     # 调用GPT获取结构化数据
#     struct_msg, _ = ask_gpt(struct_prompt, True,flag=False)
    
#     try:
#         # 增强型JSON解析
#         json_str = struct_msg.split("```json")[1].split("```")[0].strip()
#         msg_json = json.loads(json_str)
        
#         # 安全字段获取（带默认值）
#         title = msg_json.get("title", "Untitled Research")
#         authors = msg_json.get("author", [])
#         organizations = msg_json.get("organization", [])
#         abstract = msg_json.get("abstract", "暂无摘要")
#         research_background = msg_json.get("background", "暂无背景信息")
#         contribution = msg_json.get("contribution", "暂无贡献说明") 
#         methodology = msg_json.get("method", "暂无方法描述")
#         experiment = msg_json.get("experiment", "暂无实验细节")
#         conclusion = msg_json.get("conclusion", "暂无明确结论")
        
#         # 格式化处理
#         author = ", ".join(authors) if authors else "Anonymous"
#         organization = ", ".join(filter(None, organizations))  # 过滤空值
         
#     except Exception as e:
#         logging.error(f"结构化数据解析失败: {str(e)}\n原始响应：{struct_msg}")
#         return

#     # ================== 第二阶段：PPT备注生成 ==================
#     section_notes = {}
#     notes_prompt = {
#         "abstract": "用新闻播报风格总结该研究的核心摘要，突出创新价值（150字以内）",
#         "background": "以新闻导语形式阐述研究背景，说明该领域现状和研究必要性",
#         "contribution": "采用成果简报形式，用数字量化方式呈现研究贡献",
#         "method": "用技术解读风格解释方法创新点，避免专业术语",
#         "experiment": "以数据播报形式呈现实验结果，包含关键指标对比",
#         "conclusion": "用新闻总结式语气陈述结论，强调实际应用价值"    
#     }
#     for section in ["abstract", "background", "contribution", "method", "experiment", "conclusion"]:
#         prompt = f"""
#         【新闻播报备注生成 - {section}】
#         根据以下内容生成PPT备注：
#         {msg_json.get(section, '')}

#         要求：
#         1. 严格遵循：{notes_prompt[section]}
#         2. 长度控制在3-5个完整句子
#         3. 使用吸引注意力的新闻表达方式
#         4. 禁止使用项目符号或格式标记
#         """
#         note, _ = ask_gpt(prompt, True,flag=False)
#         section_notes[section] = note.strip()

    
#     #ppt_notes, _ = ask_gpt(notes_prompt, True,flag=False)

#     # ================== 图片处理模块 ==================
#     # 获取图片解释信息
#     pic, _ = ask_gpt("""取出给定markdown中所有形如“![loading][numX]”的字段及下方两行的解释：
#     1. 解释需以Fig/Figure/图开头
#     2. 输出格式：![loading][numX]`解释`位置分类
#     3. 位置分类包括：摘要、背景、贡献、方法、实验、结论
#     示例输入：
#     ![loading][num1]
#     Figure 1. 模型架构
#     示例输出：
#     ![loading][num1]`这张图展示了模型架构`方法

#     文档内容：
#     """ + content[0], True, flag=False)

#     # 分类整理图片
#     catelist = {
#         "摘要": [], "研究背景": [], "本文贡献": [], 
#         "方法": [], "实验": [], "结论": [], "其他": []
#     }
#     KEY_MAPPING = {
#         "摘要": "摘要",
#         "背景": "研究背景",
#         "贡献": "本文贡献",
#         "方法": "方法",
#         "实验": "实验",
#         "结论": "结论"
#     }

#     for line in pic.split("\n"):
#         if "![loading]" not in line:
#             continue
        
#         try:
#             # 解析格式：![loading][numX]`解释`分类
#             parts = line.split("`")
#             if len(parts) < 3:
#                 continue
                
#             raw_num = parts[0].split("[num")[1].split("]")[0]
#             picnum = int(raw_num)
#             description = parts[1]
#             category = KEY_MAPPING.get(parts[2].strip(), "其他")
            
#             # 有效性检查
#             if category not in catelist:
#                 category = "其他"
                
#             catelist[category].append((picnum, description))
            
#         except Exception as e:
#             logging.warning(f"图片解析异常：{str(e)} 行内容：{line}")
#             continue

#     # ================== PPT构建模块 ==================
#     # 创建演示文稿
#     ppt = Presentation()
#     ppt.slide_height = Inches(9)  # 16:9比例
#     ppt.slide_width = Inches(16)

#     # ------------------ 封面页设计 ------------------
#     cover_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
#     # 渐变背景
#     bg = cover_slide.background
#     fill = bg.fill
#     fill.gradient()
#     fill.gradient_angle = 45
#     stops = fill.gradient_stops
#     stops[0].color.rgb = RGBColor(12, 32, 58)   # 深蓝
#     stops[1].color.rgb = RGBColor(25, 95, 155)  # 主蓝
    
#     # 主标题
#     title_shape = cover_slide.shapes.add_textbox(
#         left=Inches(1), top=Inches(2),
#         width=Inches(14), height=Inches(3)
#     )
#     tf = title_shape.text_frame
#     tf.text = title
#     tf.paragraphs[0].font.size = Pt(44)
#     tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
#     tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
#     # 作者信息
#     subtitle_text = f"{author}\n{organization}" if organization else author
#     subtitle_shape = cover_slide.shapes.add_textbox(
#         left=Inches(4), top=Inches(5.5),
#         width=Inches(8), height=Inches(1.5)
#     )
#     stf = subtitle_shape.text_frame
#     stf.text = subtitle_text
#     stf.paragraphs[0].font.size = Pt(24)
#     stf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
#     stf.paragraphs[0].alignment = PP_ALIGN.CENTER

#     # ------------------ 内容页模板 ------------------
#     sections = [
#         ("摘要", abstract, catelist["摘要"], section_notes["abstract"]),
#         ("研究背景", research_background, catelist["研究背景"], section_notes["background"]),
#         ("本文贡献", contribution, catelist["本文贡献"], section_notes["contribution"]),
#         ("方法设计", methodology, catelist["方法"], section_notes["method"]),
#         ("实验验证", experiment, catelist["实验"], section_notes["experiment"]),
#         ("研究结论", conclusion, catelist["结论"], section_notes["conclusion"])
#     ]

#     for section_title, content_text, images, note_text in sections:
#         # 创建内容页
#         content_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
#         notes_slide = content_slide.notes_slide
#         text_frame = notes_slide.notes_text_frame
#         text_frame.text = note_text  # 关联当前页备注
#         # 统一背景
#         content_slide.background.fill.solid()
#         content_slide.background.fill.fore_color.rgb = RGBColor(240, 245, 249)  # 浅灰蓝
        
#         # 添加备注
#         if note_text:
#             content_slide.notes_slide.notes_text_frame.text = note_text
            
#         # 内容区域
#         text_box = content_slide.shapes.add_textbox(
#             left=Inches(1), top=Inches(1.5),
#             width=Inches(10), height=Inches(5))
#         tf = text_box.text_frame
#         tf.word_wrap = True
#         tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
#         # 段落格式
#         for para in content_text.split("\n"):
#             p = tf.add_paragraph()
#             p.text = para.strip()
#             p.font.size = Pt(22)
#             p.font.color.rgb = RGBColor(51, 51, 51)
#             p.line_spacing = 1.3
            
#         # 插入相关图片
#         for idx, (picnum, desc) in enumerate(images):
#             try:
#                 img_data = pics[picnum-1].split(",", 1)[1]
#                 img_bytes = base64.b64decode(img_data)
                
#                 # 创建图片幻灯片
#                 img_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
                
#                 # 添加图片
#                 img_stream = io.BytesIO(img_bytes)
#                 img_shape = img_slide.shapes.add_picture(
#                     img_stream, 
#                     left=Inches(1), top=Inches(1),
#                     width=Inches(12), height=Inches(6.75)  # 保持16:9比例
#                 )
                
#                 # 图片边框
#                 img_shape.line.color.rgb = RGBColor(220, 220, 220)
#                 img_shape.line.width = Pt(1.5)
                
#                 # 图片说明
#                 caption_box = img_slide.shapes.add_textbox(
#                     left=Inches(2), top=Inches(7.5),
#                     width=Inches(12), height=Inches(1)
#                 )
#                 caption_box.text = desc
#                 caption_box.text_frame.paragraphs[0].font.size = Pt(18)
#                 caption_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
                
#             except Exception as e:
#                 logging.error(f"图片插入失败：{str(e)}")
#                 continue

#     # ================== 最终保存 ==================
#     output_path = os.path.join(dir_ppt, f"{Path(file_name).stem}_presentation.pptx")
#     ppt.save(output_path)
#     print(f"\n✅ PPT生成成功！保存路径：{output_path}")    
def wrap_text(text, max_chars_per_line=20):
    import textwrap
    return "\n".join(textwrap.wrap(text, width=max_chars_per_line))
 
def generate_ppt(content, file_name, dir_ppt):
    """生成研究论文PPT的完整函数
    
    Args:
        content (list): 包含文档内容和图片的列表
        file_name (str): 原始文件名
        dir_ppt (str): PPT输出目录
    """
    # ================== 初始化处理 ==================
    # 提取图片数据
    pics = ("[num1]:data:image" + content[1]).replace("\n\n", "\n").split("\n")

    # ================== 第一阶段：结构化数据提取 ==================
    struct_prompt = f"""
    【严格JSON生成任务】
    请从以下文档中提取：
    - title（英文原题）
    - author（作者列表）
    - organization（单位列表，无则保留空字符串）
    - abstract（中文摘要）
    - background（研究背景）
    - contribution（本文贡献）
    - method（方法）
    - experiment（实验）
    - conclusion（结论）

    输出要求：
    1. 严格遵循JSON格式
    2. 使用```json代码块包裹
    3. 除标题作者单位外全用中文
    4. 空值字段保持空字符串

    文档内容：
    {content[0]}
    """
    
    # 调用GPT获取结构化数据
    struct_msg, _ = ask_gpt(struct_prompt, True,flag=False)
    
    try:
        # 增强型JSON解析
        json_str = struct_msg.split("```json")[1].split("```")[0].strip()
        msg_json = json.loads(json_str)
        
        # 安全字段获取（带默认值）
        title = msg_json.get("title", "Untitled Research")
        authors = msg_json.get("author", [])
        organizations = msg_json.get("organization", [])
        abstract = msg_json.get("abstract", "暂无摘要")
        research_background = msg_json.get("background", "暂无背景信息")
        contribution = msg_json.get("contribution", "暂无贡献说明") 
        methodology = msg_json.get("method", "暂无方法描述")
        experiment = msg_json.get("experiment", "暂无实验细节")
        conclusion = msg_json.get("conclusion", "暂无明确结论")
        
        # 格式化处理
        author = ", ".join(authors) if authors else "Anonymous"
        #organization = ", ".join(filter(None, organizations))  # 过滤空值
        organization = ", ".join(sorted(set(filter(None, organizations))))
         
    except Exception as e:
        logging.error(f"结构化数据解析失败: {str(e)}\n原始响应：{struct_msg}")
        return

    # ================== 第二阶段：PPT备注生成 ==================
    section_notes = {}
    notes_prompt = {
        "abstract": "用新闻播报风格总结该研究的核心摘要，突出创新价值（150字以内）",
        "background": "以新闻导语形式阐述研究背景，说明该领域现状和研究必要性",
        "contribution": "采用成果简报形式，用数字量化方式呈现研究贡献",
        "method": "用技术解读风格解释方法创新点，避免专业术语",
        "experiment": "以数据播报形式呈现实验结果，包含关键指标对比",
        "conclusion": "用新闻总结式语气陈述结论，强调实际应用价值"    
    }
    for section in ["abstract", "background", "contribution", "method", "experiment", "conclusion"]:
        prompt = f"""
        【新闻播报备注生成 - {section}】
        根据以下内容生成PPT备注：
        {msg_json.get(section, '')}

        要求：
        1. 严格遵循：{notes_prompt[section]}
        2. 长度控制在3-5个完整句子
        3. 使用吸引注意力的新闻表达方式
        4. 禁止使用项目符号或格式标记
        """
        note, _ = ask_gpt(prompt, True,flag=False)
        section_notes[section] = note.strip()

    
    #ppt_notes, _ = ask_gpt(notes_prompt, True,flag=False)

    # ================== 图片处理模块 ==================
    # 获取图片解释信息
    pic, _ = ask_gpt("""取出给定markdown中所有形如“![loading][numX]”的字段及下方两行的解释：
    1. 解释需以Fig/Figure/图开头
    2. 输出格式：![loading][numX]`解释`位置分类
    3. 位置分类包括：摘要、背景、贡献、方法、实验、结论
    示例输入：
    ![loading][num1]
    Figure 1. 模型架构
    示例输出：
    ![loading][num1]`这张图展示了模型架构`方法

    文档内容：
    """ + content[0], True, flag=False)

    # 分类整理图片
    catelist = {
        "摘要": [], "研究背景": [], "本文贡献": [], 
        "方法": [], "实验": [], "结论": [], "其他": []
    }
    KEY_MAPPING = {
        "摘要": "摘要",
        "背景": "研究背景",
        "贡献": "本文贡献",
        "方法": "方法",
        "实验": "实验",
        "结论": "结论"
    }

    for line in pic.split("\n"):
        if "![loading]" not in line:
            continue
        
        try:
            # 解析格式：![loading][numX]`解释`分类
            parts = line.split("`")
            if len(parts) < 3:
                continue
                
            raw_num = parts[0].split("[num")[1].split("]")[0]
            picnum = int(raw_num)
            description = parts[1]
            category = KEY_MAPPING.get(parts[2].strip(), "其他")
            
            # 有效性检查
            if category not in catelist:
                category = "其他"
                
            catelist[category].append((picnum, description))
            
        except Exception as e:
            logging.warning(f"图片解析异常：{str(e)} 行内容：{line}")
            continue

    # ================== PPT构建模块 ==================
    # 创建演示文稿
    ppt = Presentation()
    ppt.slide_height = Inches(9)  # 16:9比例
    ppt.slide_width = Inches(16)

    # ------------------ 封面页设计 ------------------
    cover_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # 渐变背景
#     bg = cover_slide.background
#     fill = bg.fill
#     fill.gradient()
#     fill.gradient_angle = 45
#     stops = fill.gradient_stops
#     stops[0].color.rgb = RGBColor(12, 32, 58)   # 深蓝
#     stops[1].color.rgb = RGBColor(25, 95, 155)  # 主蓝
    
    # 主标题
    title_box = cover_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(12), Inches(3))
    text_frame = title_box.text_frame
    text_frame.clear()
  # 清除默认段落
    if len(title) > 30:
        midpoint = len(title) // 2
        split_index = title[:midpoint].rfind(' ') + 1
        if split_index == 0:
            split_index = midpoint
        line1 = title[:split_index]
        line2 = title[split_index:]
        text_frame.text = f"{line1}\n{line2}"
    else:
        text_frame.text = title

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(45)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.alignment = PP_ALIGN.CENTER
        
        
        
        
    # 作者信息
    subtitle_text = f"{author}\n{organization}" if organization else author
    subtitle_shape = cover_slide.shapes.add_textbox(Inches(4), Inches(5.5),Inches(8), Inches(1.5))
    stf = subtitle_shape.text_frame
    stf.text = subtitle_text
    stf.paragraphs[0].font.size = Pt(24)
    stf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ------------------ 内容页模板 ------------------
    sections = [
        ("摘要", abstract, catelist["摘要"], section_notes["abstract"]),
        ("研究背景", research_background, catelist["研究背景"], section_notes["background"]),
        ("本文贡献", contribution, catelist["本文贡献"], section_notes["contribution"]),
        ("方法设计", methodology, catelist["方法"], section_notes["method"]),
        ("实验验证", experiment, catelist["实验"], section_notes["experiment"]),
        ("研究结论", conclusion, catelist["结论"], section_notes["conclusion"])
    ]

    for section_title, content_text, images, note_text in sections:
        # 创建内容页
        content_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        notes_slide = content_slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = note_text  # 关联当前页备注
        # 统一背景
        content_slide.background.fill.solid()
        content_slide.background.fill.fore_color.rgb = RGBColor(240, 245, 249)  # 浅灰蓝
        
        # 添加备注
        if note_text:
            content_slide.notes_slide.notes_text_frame.text = note_text
            
        # 内容区域
        text_box = content_slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1.5),
            width=Inches(10), height=Inches(5))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # 段落格式
        for para in content_text.split("\n"):
            p = tf.add_paragraph()
            p.text = para.strip()
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.line_spacing = 1.3
            
        # 插入相关图片
        for idx, (picnum, desc) in enumerate(images):
            try:
                img_data = pics[picnum-1].split(",", 1)[1]
                img_bytes = base64.b64decode(img_data)
                
                # 创建图片幻灯片
                img_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
                
                # 添加图片
                img_stream = io.BytesIO(img_bytes)
                img_shape = img_slide.shapes.add_picture(
                    img_stream, 
                    left=Inches(1), top=Inches(1),
                    width=Inches(12), height=Inches(6.75)  # 保持16:9比例
                )
                
                # 图片边框
                img_shape.line.color.rgb = RGBColor(220, 220, 220)
                img_shape.line.width = Pt(1.5)
                
                
                # 图片说明
                wrapped_desc = wrap_text(desc, max_chars_per_line=30)
                caption_box = img_slide.shapes.add_textbox(
                    left=Inches(2), top=Inches(8),
                    width=Inches(12), height=Inches(1)
                )
                tf = caption_box.text_frame
                tf.text = wrapped_desc
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(18)
                    paragraph.font.color.rgb = RGBColor(100, 100, 100)
                    paragraph.alignment = PP_ALIGN.CENTER
#                 caption_box.text = desc
#                 caption_box.text_frame.paragraphs[0].font.size = Pt(18)
#                 caption_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
                
                tf = caption_box.text_frame
                tf.clear()  # 清除原内容，避免重复
                tf.word_wrap = True  # 启用自动换行
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 内容自动适应文本框
                p = tf.add_paragraph()
                p.text = desc
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(100, 100, 100)
                p.alignment = PP_ALIGN.CENTER
                
                
                
            except Exception as e:
                logging.error(f"图片插入失败：{str(e)}")
                continue

    # ================== 最终保存 ==================
    output_path = os.path.join(dir_ppt, f"{Path(file_name).stem}_presentation.pptx")
    ppt.save(output_path)
    print(f"\n✅ PPT生成成功！保存路径：{output_path}")    
    


## =============================================================================== pdf2md
## PDF转markdown
def convert_pdf_to_md(file_path, dir_md):
    md_name = file_path.split("/")[-1].replace(".pdf", ".md").replace(" ", "_").replace("(", "").replace(")", "")
    command = [
        'curl',
        '-X', 'POST',
        'http://127.0.0.1:8000/upload-pdf/',
        '-F', f'file=@{file_path}',
        '-F', f'output_file={dir_md}/'
    ]
    while True:
        result = subprocess.run(command, capture_output=True, text=True)
        return_code = result.returncode
        if return_code == 0:
            print(f"convert pdf to md: {os.path.join(dir_md, md_name)}")
            newapi_pdf2md_process(dir_md + "/" + md_name.replace(".md", ""), dir_md)
            return dir_md + "/" + md_name, md_name
        
        


## 处理新API的文件夹，把md修改了一下存在md路径下
def newapi_pdf2md_process(folder_path, md_dir_path):
    title = folder_path.split("/")[-1]
    newapi_mdpath = os.path.join(folder_path, "auto", title + ".md")
    with open(newapi_mdpath, 'r', encoding='utf-8') as file:
        content = file.read()
        
    def replace_images(text):
        counter = [1]
        original_paths = []

        def replace(match):
            original_image = match.group(1)
            original_paths.append(original_image)
            result = f'![loading][num{counter[0]}]'
            counter[0] += 1
            return result

        pattern = r'!\[\]\((images/[^)]+\.jpg)\)'
        modified_text = re.sub(pattern, replace, text)

        return modified_text, original_paths
    
    new_content, original_path = replace_images(content)
    imgfolder_path = [os.path.join(folder_path, "auto", i) for i in original_path]
    
    def image_to_base64(image_path):
        with open(image_path, 'rb') as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
        return encoded_string
    
    base64string = ""
    for i in range(len(imgfolder_path)):
        base64string += f"[num{i+1}]:data:image/jpeg;base64," + image_to_base64(imgfolder_path[i]) + "\n\n"
    
    final_mdpath = os.path.join(md_dir_path, title + ".md")
    with open(final_mdpath, 'w', encoding='utf-8') as file:
        file.write(new_content + "\n\n\n" + base64string)
        
    try:
        subprocess.run(['rm', '-rf', folder_path], check=True)
        print(f"目录 '{folder_path}' 已成功删除。")
    except subprocess.CalledProcessError as e:
        print(f"发生错误：{e}")
    except Exception as e:
        print(f"发生其他错误：{e}")
        


## =============================================================================== workflow

def workflow(file_name, dir_md, dir_ppt, pdf2md=False):
    print("\n\n\n==============" + file_name.name + "==============")
    if pdf2md:
        md_path, md_name = convert_pdf_to_md(str(file_name), dir_md)
    with open(md_path, 'r') as file:
        content = file.read().split("[num1]:data:image")

    generate_ppt(content, md_name, dir_ppt)
    time.sleep(1)


## 判断文件夹是否存在，不存在则创建文件夹
def make_dir(folder):
    if not os.path.exists(folder):
        os.makedirs(folder)
    

if __name__ == '__main__':
    dir_pdf = "/home/asus/HuFan/ppt/pdf"
    dir_md = "/home/asus/HuFan/ppt/md"
    dir_ppt = "/home/asus/HuFan/ppt/ppt1"

    make_dir(dir_md)
    make_dir(dir_ppt)
    
    folder = Path(dir_pdf)
    pdf_files = folder.glob('*.pdf')
    
    for pdf_file in pdf_files:
        workflow(pdf_file, dir_md, dir_ppt, pdf2md=True)